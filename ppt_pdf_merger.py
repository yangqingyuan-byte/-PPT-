import os
import sys
import subprocess
import datetime
import tempfile
import shutil
import time
import json
import platform
from dataclasses import dataclass
from typing import List, Optional, Tuple

import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk

try:
    import ttkbootstrap as ttkb
except ImportError:  # pragma: no cover
    ttkb = None

try:
    import PyPDF2
except ImportError:  # pragma: no cover
    PyPDF2 = None

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
except ImportError:  # pragma: no cover
    A4 = None
    canvas = None
    pdfmetrics = None
    TTFont = None

try:
    import win32com.client
except ImportError:  # pragma: no cover
    win32com = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None


@dataclass
class PPTItem:
    display_name: str
    file_path: str


class DraggableListbox(tk.Listbox):
    def __init__(self, master, **kwargs):
        super().__init__(master, **kwargs)
        self._dragging = False
        self._drag_start_index: Optional[int] = None
        self.bind("<ButtonPress-1>", self._on_button_press)
        self.bind("<ButtonRelease-1>", self._on_button_release)
        self.bind("<B1-Motion>", self._on_motion)

    def _on_button_press(self, event):
        self._dragging = True
        self._drag_start_index = self.nearest(event.y)

    def _on_motion(self, event):
        if not self._dragging or self._drag_start_index is None:
            return
        new_index = self.nearest(event.y)
        if new_index == self._drag_start_index or new_index < 0:
            return
        item_text = self.get(self._drag_start_index)
        self.delete(self._drag_start_index)
        self.insert(new_index, item_text)
        self.selection_clear(0, tk.END)
        self.selection_set(new_index)
        self._drag_start_index = new_index
        self.event_generate("<<ListboxReordered>>")

    def _on_button_release(self, _event):
        self._dragging = False
        self._drag_start_index = None


class PPTMergerApp:
    def __init__(self, root: tk.Tk):
        self.root = root
        self.root.title("PPT 转 PDF 合并工具")
        self.root.minsize(780, 460)

        ttkb_window_cls = getattr(ttkb, "Window", None)
        self.use_bootstrap = ttkb_window_cls is not None and isinstance(self.root, ttkb_window_cls)
        self.style = ttk.Style(self.root)
        self._style_map = {
            "success": "Success.TButton",
            "primary": "Primary.TButton",
            "info": "Info.TButton",
            "secondary": "Secondary.TButton",
        }
        self._configure_styles()

        self.script_dir = os.path.abspath(os.path.dirname(sys.argv[0] or __file__))
        self.vbs_path = os.path.normpath(os.path.join(self.script_dir, "单个ppt转为pdf.vbs"))
        self.config_path = os.path.join(self.script_dir, "ppt_merger_settings.json")
        self.font_regular = "Helvetica"
        self.font_bold = "Helvetica-Bold"
        self._font_checked = False
        self.is_windows = platform.system() == "Windows"
        self.is_mac = platform.system() == "Darwin"

        self.folder_path: Optional[str] = None
        self.available_items: List[PPTItem] = []
        self.selected_items: List[PPTItem] = []

        self._build_ui()
        self._ensure_chinese_font()
        self._load_last_state()

    def _build_ui(self):
        outer = ttk.Frame(self.root, padding=(12, 12))
        outer.pack(fill=tk.BOTH, expand=True)

        chooser_frame = ttk.Frame(outer)
        chooser_frame.pack(fill=tk.X)

        ttk.Label(chooser_frame, text="当前目录：").pack(side=tk.LEFT)
        self.folder_var = tk.StringVar(value="尚未选择")
        folder_entry = ttk.Entry(chooser_frame, textvariable=self.folder_var, state="readonly")
        folder_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(4, 8))

        self._create_button(chooser_frame, text="选择目录", command=self.choose_folder, bootstyle="primary").pack(
            side=tk.LEFT
        )

        lists_frame = ttk.Frame(outer)
        lists_frame.pack(fill=tk.BOTH, expand=True, pady=12)

        # 可选列表
        left_frame = ttk.Frame(lists_frame)
        left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        ttk.Label(left_frame, text="可选 PPT 文件").pack()
        self.available_listbox = tk.Listbox(left_frame, selectmode=tk.EXTENDED)
        self.available_listbox.pack(fill=tk.BOTH, expand=True, padx=6, pady=6)

        # 中间按钮
        middle_frame = ttk.Frame(lists_frame)
        middle_frame.pack(side=tk.LEFT, fill=tk.Y, padx=10)
        self._create_button(middle_frame, text="全选 →", command=self.add_all, bootstyle="secondary").pack(
            pady=6, fill=tk.X
        )
        self._create_button(middle_frame, text="添加 →", command=self.add_selected, bootstyle="primary").pack(
            pady=6, fill=tk.X
        )
        self._create_button(middle_frame, text="← 移除", command=self.remove_selected, bootstyle="info").pack(
            pady=6, fill=tk.X
        )
        self._create_button(middle_frame, text="清空", command=self.clear_selected, bootstyle="secondary").pack(
            pady=6, fill=tk.X
        )

        # 已选列表
        right_frame = ttk.Frame(lists_frame)
        right_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        ttk.Label(right_frame, text="已选 PPT 文件（可拖拽排序）").pack()
        self.selected_listbox = DraggableListbox(right_frame, selectmode=tk.BROWSE)
        self.selected_listbox.pack(fill=tk.BOTH, expand=True, padx=6, pady=6)
        self.selected_listbox.bind("<<ListboxReordered>>", self._sync_order_with_model)

        # 底部按钮
        bottom_frame = ttk.Frame(outer)
        bottom_frame.pack(fill=tk.X)

        # 第一行：合并PPT按钮
        merge_ppt_frame = ttk.Frame(bottom_frame)
        merge_ppt_frame.pack(fill=tk.X, pady=(0, 4))
        
        self._create_button(
            merge_ppt_frame,
            text="合并为 PPT",
            command=self.merge_ppts,
            bootstyle="info",
        ).pack(side=tk.LEFT, fill=tk.X, expand=True, padx=4)

        # 第二行：PDF合并按钮
        pdf_frame = ttk.Frame(bottom_frame)
        pdf_frame.pack(fill=tk.X)

        self._create_button(
            pdf_frame,
            text="博士组会",
            command=lambda: self.start_process("博士组会"),
            bootstyle="success",
        ).pack(side=tk.LEFT, fill=tk.X, expand=True, padx=4, pady=(0, 4))

        self._create_button(
            pdf_frame,
            text="大模型和开放世界组组会",
            command=lambda: self.start_process("大模型和开放世界组组会"),
            bootstyle="primary",
        ).pack(side=tk.LEFT, fill=tk.X, expand=True, padx=4, pady=(0, 4))

    def _create_button(self, parent, text, command, bootstyle="secondary", **kwargs):
        if self.use_bootstrap and ttkb is not None:
            return ttkb.Button(parent, text=text, command=command, bootstyle=bootstyle, **kwargs)

        style_name = self._style_map.get(bootstyle.lower(), "TButton")
        button = ttk.Button(parent, text=text, command=command, style=style_name, **kwargs)
        return button

    def _configure_styles(self):
        if self.use_bootstrap:
            return

        try:
            self.style.theme_use("clam")
        except tk.TclError:
            pass

        default_font = ("Microsoft YaHei", 11)
        for style_name in {"TButton", *self._style_map.values()}:
            self.style.configure(style_name, font=default_font, padding=6)

        self.style.configure("Success.TButton", background="#4CAF50", foreground="white")
        self.style.map(
            "Success.TButton",
            background=[("pressed", "#388E3C"), ("active", "#45A049")],
            foreground=[("disabled", "#DDDDDD")],
        )

        self.style.configure("Primary.TButton", background="#2196F3", foreground="white")
        self.style.map(
            "Primary.TButton",
            background=[("pressed", "#1976D2"), ("active", "#1E88E5")],
            foreground=[("disabled", "#DDDDDD")],
        )

        self.style.configure("Info.TButton", background="#00ACC1", foreground="white")
        self.style.map(
            "Info.TButton",
            background=[("pressed", "#00838F"), ("active", "#0097A7")],
            foreground=[("disabled", "#DDDDDD")],
        )

        self.style.configure("Secondary.TButton", background="#607D8B", foreground="white")
        self.style.map(
            "Secondary.TButton",
            background=[("pressed", "#455A64"), ("active", "#546E7A")],
            foreground=[("disabled", "#DDDDDD")],
        )

    def choose_folder(self):
        folder = filedialog.askdirectory()
        if not folder:
            return
        self.folder_path = folder
        self.folder_var.set(folder)
        self._load_ppt_files()
        self._save_last_state()

    def _load_last_state(self):
        if not os.path.exists(self.config_path):
            return
        try:
            with open(self.config_path, "r", encoding="utf-8") as cfg_file:
                data = json.load(cfg_file)
        except (OSError, json.JSONDecodeError):
            return
        folder = data.get("last_folder")
        if not folder or not os.path.isdir(folder):
            return
        self.folder_path = folder
        self.folder_var.set(folder)
        self._load_ppt_files()

    def _save_last_state(self):
        data = {"last_folder": self.folder_path}
        try:
            with open(self.config_path, "w", encoding="utf-8") as cfg_file:
                json.dump(data, cfg_file, ensure_ascii=False, indent=2)
        except OSError:
            pass

    def _load_ppt_files(self):
        if not self.folder_path:
            return
        self.available_items.clear()
        self.available_listbox.delete(0, tk.END)
        self.clear_selected()

        for entry in sorted(os.listdir(self.folder_path)):
            if entry.lower().endswith((".ppt", ".pptx")):
                full_path = os.path.join(self.folder_path, entry)
                if os.path.isfile(full_path):
                    item = PPTItem(display_name=entry, file_path=full_path)
                    self.available_items.append(item)
                    self.available_listbox.insert(tk.END, entry)

        if not self.available_items:
            messagebox.showinfo("提示", "该目录中未找到 PPT 或 PPTX 文件。")

    def add_selected(self):
        indices = list(self.available_listbox.curselection())
        if not indices:
            messagebox.showwarning("提示", "请在左侧列表中选择至少一个 PPT。")
            return

        for idx in indices:
            item = self.available_items[idx]
            if item not in self.selected_items:
                self.selected_items.append(item)
                self.selected_listbox.insert(tk.END, item.display_name)

    def add_all(self):
        if not self.available_items:
            messagebox.showinfo("提示", "当前目录没有可用的 PPT。")
            return
        added = False
        for item in self.available_items:
            if item not in self.selected_items:
                self.selected_items.append(item)
                self.selected_listbox.insert(tk.END, item.display_name)
                added = True
        if not added:
            messagebox.showinfo("提示", "所有 PPT 已经在右侧列表中。")

    def remove_selected(self):
        idx = self.selected_listbox.curselection()
        if not idx:
            messagebox.showwarning("提示", "请在右侧列表中选择要移除的 PPT。")
            return
        pos = idx[0]
        self.selected_listbox.delete(pos)
        del self.selected_items[pos]

    def clear_selected(self):
        self.selected_listbox.delete(0, tk.END)
        self.selected_items.clear()

    def _ensure_chinese_font(self):
        if self._font_checked:
            return
        self._font_checked = True

        if canvas is None or pdfmetrics is None or TTFont is None:
            return

        font_dir = os.path.join(os.environ.get("WINDIR", r"C:\Windows"), "Fonts")
        candidates = [
            ("SimSun", "simsun.ttc"),
            ("SimHei", "simhei.ttf"),
            ("FangSong", "simfang.ttf"),
            ("KaiTi", "simkai.ttf"),
            ("MicrosoftYaHei", "msyh.ttc"),
        ]

        for font_name, file_name in candidates:
            font_path = os.path.join(font_dir, file_name)
            if not os.path.exists(font_path):
                continue
            try:
                pdfmetrics.registerFont(TTFont(font_name, font_path))
                self.font_regular = font_name
                self.font_bold = font_name
                return
            except Exception:
                continue

    def _sync_order_with_model(self, _event=None):
        new_order: List[PPTItem] = []
        for i in range(self.selected_listbox.size()):
            name = self.selected_listbox.get(i)
            match = next((item for item in self.selected_items if item.display_name == name), None)
            if match:
                new_order.append(match)
        # 如果拖动后有重复或遗漏，回退到线性搜索结果
        if len(new_order) == len(self.selected_items):
            self.selected_items = new_order

    def start_process(self, mode_label: str):
        if not self.selected_items:
            messagebox.showwarning("提示", "请先选择至少一个 PPT 文件。")
            return

        if not self.folder_path:
            messagebox.showwarning("提示", "请先选择工作目录。")
            return

        if not os.path.exists(self.vbs_path):
            messagebox.showerror("错误", f"未找到 VBS 脚本：{self.vbs_path}")
            return

        if PyPDF2 is None or canvas is None or A4 is None:
            messagebox.showerror(
                "缺少依赖",
                "请先安装依赖库：\n\npip install PyPDF2 reportlab",
            )
            return

        try:
            stats = self._convert_ppts_to_pdfs()
            if not stats:
                return
            output_path = self._merge_pdfs_with_toc(stats, mode_label)
        except Exception as exc:  # pragma: no cover - GUI error display
            messagebox.showerror("错误", f"处理过程中出现问题：\n{exc}")
            return

        messagebox.showinfo("完成", f"合并文件已生成：\n{output_path}")

    def merge_ppts(self):
        """使用PowerPoint COM接口直接合并PPT文件"""
        if not self.selected_items:
            messagebox.showwarning("提示", "请先选择至少一个 PPT 文件。")
            return

        if not self.folder_path:
            messagebox.showwarning("提示", "请先选择工作目录。")
            return

        # 检查平台和依赖
        if self.is_windows:
            if win32com is None:
                messagebox.showerror(
                    "缺少依赖",
                    "请先安装依赖库：\n\npip install pywin32",
                )
                return
        elif self.is_mac:
            # Mac 上使用 python-pptx
            if Presentation is None:
                messagebox.showerror(
                    "缺少依赖",
                    "请先安装依赖库：\n\npip install python-pptx",
                )
                return
        else:
            messagebox.showerror(
                "不支持的操作系统",
                f"当前操作系统 {platform.system()} 暂不支持合并PPT功能。\n请使用 Windows 或 macOS。",
            )
            return

        try:
            output_path = self._merge_ppts_with_com()
            messagebox.showinfo("完成", f"合并PPT文件已生成：\n{output_path}")
        except Exception as exc:
            messagebox.showerror("错误", f"合并PPT过程中出现问题：\n{exc}")

    def _merge_ppts_with_com(self) -> str:
        """合并PPT文件（Windows使用COM接口，Mac使用python-pptx）"""
        if self.is_windows:
            return self._merge_ppts_windows()
        elif self.is_mac:
            return self._merge_ppts_mac()
        else:
            raise RuntimeError(f"不支持的操作系统: {platform.system()}")

    def _merge_ppts_windows(self) -> str:
        """使用PowerPoint COM接口合并PPT文件（Windows）"""
        today_str = datetime.datetime.now().strftime("%Y%m%d")
        output_name = f"{today_str}合并PPT.pptx"
        output_path = os.path.join(self.folder_path, output_name)

        # 如果文件已存在，添加序号
        counter = 1
        base_output_path = output_path
        while os.path.exists(output_path):
            name_without_ext = os.path.splitext(base_output_path)[0]
            output_path = f"{name_without_ext}_{counter}.pptx"
            counter += 1

        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        # 尝试隐藏窗口，如果失败则忽略（某些版本的PowerPoint不允许隐藏）
        try:
            ppt_app.Visible = False
        except Exception:
            pass  # 如果无法隐藏窗口，继续执行（窗口会显示）

        try:
            # 打开第一个PPT作为主文件
            first_item = self.selected_items[0]
            first_path = os.path.abspath(os.path.normpath(first_item.file_path))
            main_presentation = ppt_app.Presentations.Open(first_path, WithWindow=False)

            # 统计信息：用于创建目录页
            slide_counts = []
            slide_counts.append((first_item.display_name, main_presentation.Slides.Count))

            # 复制其他PPT的幻灯片
            for item in self.selected_items[1:]:
                ppt_path = os.path.abspath(os.path.normpath(item.file_path))
                source_presentation = ppt_app.Presentations.Open(ppt_path, WithWindow=False)
                
                slide_count = source_presentation.Slides.Count
                slide_counts.append((item.display_name, slide_count))

                # 复制所有幻灯片到主文件
                for i in range(1, slide_count + 1):
                    source_slide = source_presentation.Slides(i)
                    source_slide.Copy()
                    # 粘贴到主文件末尾
                    main_presentation.Slides.Paste()
                    # 保持原幻灯片的布局和格式
                    pasted_slide = main_presentation.Slides(main_presentation.Slides.Count)
                    try:
                        pasted_slide.Design = source_slide.Design
                    except Exception:
                        pass  # 某些设计可能无法复制，忽略错误
                    try:
                        pasted_slide.ColorScheme = source_slide.ColorScheme
                    except Exception:
                        pass  # 某些配色方案可能无法复制，忽略错误

                source_presentation.Close()

            # 创建目录页（插入到第一页）
            self._create_toc_slide(main_presentation, slide_counts)

            # 保存合并后的PPT
            main_presentation.SaveAs(output_path)
            main_presentation.Close()

            return output_path
        finally:
            ppt_app.Quit()

    def _create_toc_slide(self, presentation, slide_counts: List[Tuple[str, int]]):
        """在PPT中创建目录页"""
        try:
            # 在开头插入新幻灯片（使用空白布局）
            toc_slide = presentation.Slides.Add(1, 5)  # 5 = ppLayoutBlank
            
            # 添加标题文本框
            title_left = 72
            title_top = 72
            title_width = presentation.PageSetup.SlideWidth - 144
            title_height = 80
            
            title_box = toc_slide.Shapes.AddTextbox(1, title_left, title_top, title_width, title_height)
            title_range = title_box.TextFrame.TextRange
            title_range.Text = "目录"
            title_range.Font.Size = 44
            title_range.Font.Bold = True
            title_range.Font.Name = "Microsoft YaHei"

            # 添加内容文本框
            content_left = 72
            content_top = 180
            content_width = presentation.PageSetup.SlideWidth - 144
            content_height = presentation.PageSetup.SlideHeight - 250

            text_box = toc_slide.Shapes.AddTextbox(1, content_left, content_top, content_width, content_height)
            text_frame = text_box.TextFrame
            text_frame.WordWrap = 1  # 自动换行
            text_frame.AutoSize = 0  # 不自动调整大小

            # 构建目录内容
            cumulative_slide = 2  # 从第2页开始（第1页是目录页）
            toc_lines = []
            
            for idx, (display_name, slide_count) in enumerate(slide_counts, start=1):
                start_slide = cumulative_slide
                toc_lines.append(f"{idx}. {display_name}  页数: {slide_count}  起始页: {start_slide}")
                cumulative_slide += slide_count

            # 设置文本内容
            text_range = text_frame.TextRange
            text_range.Text = "\n".join(toc_lines)
            
            # 设置字体大小和格式
            text_range.Font.Size = 24
            text_range.Font.Name = "Microsoft YaHei"
            text_range.ParagraphFormat.SpaceAfter = 6  # 段落间距

            # 设置行距
            try:
                for i in range(1, len(toc_lines) + 1):
                    para = text_range.Paragraphs(i)
                    para.ParagraphFormat.LineSpacing = 28  # 行距
            except Exception:
                # 如果设置行距失败，使用默认值
                pass

        except Exception as e:
            # 如果创建目录页失败，不影响主流程，只记录错误
            print(f"创建目录页时出错：{e}")

    def _merge_ppts_mac(self) -> str:
        """使用python-pptx合并PPT文件（Mac）"""
        today_str = datetime.datetime.now().strftime("%Y%m%d")
        output_name = f"{today_str}合并PPT.pptx"
        output_path = os.path.join(self.folder_path, output_name)

        # 如果文件已存在，添加序号
        counter = 1
        base_output_path = output_path
        while os.path.exists(output_path):
            name_without_ext = os.path.splitext(base_output_path)[0]
            output_path = f"{name_without_ext}_{counter}.pptx"
            counter += 1

        # 打开第一个PPT作为主文件
        first_item = self.selected_items[0]
        first_path = os.path.abspath(os.path.normpath(first_item.file_path))
        main_presentation = Presentation(first_path)

        # 统计信息：用于创建目录页
        slide_counts = []
        slide_counts.append((first_item.display_name, len(main_presentation.slides)))

        # 复制其他PPT的幻灯片
        for item in self.selected_items[1:]:
            ppt_path = os.path.abspath(os.path.normpath(item.file_path))
            source_presentation = Presentation(ppt_path)
            
            slide_count = len(source_presentation.slides)
            slide_counts.append((item.display_name, slide_count))

            # 复制所有幻灯片到主文件
            # 注意：python-pptx的复制功能有限，这里使用XML直接复制的方式
            for source_slide in source_presentation.slides:
                # 直接复制整个幻灯片的XML元素
                import copy
                from lxml import etree
                
                # 创建新幻灯片，使用源幻灯片的布局
                slide_layout = source_slide.slide_layout
                new_slide = main_presentation.slides.add_slide(slide_layout)
                
                # 复制源幻灯片的XML内容
                source_xml = source_slide.element
                new_xml = copy.deepcopy(source_xml)
                
                # 替换新幻灯片的XML
                new_slide.element.getparent().replace(new_slide.element, new_xml)
                new_slide.element = new_xml

        # 创建目录页（插入到第一页）
        self._create_toc_slide_pptx(main_presentation, slide_counts)

        # 保存合并后的PPT
        main_presentation.save(output_path)
        return output_path

    def _create_toc_slide_pptx(self, presentation, slide_counts: List[Tuple[str, int]]):
        """在PPT中创建目录页（使用python-pptx）"""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            # 获取空白布局
            blank_layout = presentation.slide_layouts[6]  # 6 = 空白布局
            toc_slide = presentation.slides.add_slide(blank_layout)

            # 添加标题
            left = Inches(1)
            top = Inches(1)
            width = Inches(8)
            height = Inches(0.8)
            
            title_box = toc_slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = "目录"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.name = "Microsoft YaHei"
            title_para.alignment = PP_ALIGN.LEFT

            # 添加内容文本框
            content_left = Inches(1)
            content_top = Inches(2)
            content_width = Inches(8)
            content_height = Inches(5)

            text_box = toc_slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            # 构建目录内容
            cumulative_slide = 2  # 从第2页开始（第1页是目录页）
            toc_lines = []
            
            for idx, (display_name, slide_count) in enumerate(slide_counts, start=1):
                start_slide = cumulative_slide
                toc_lines.append(f"{idx}. {display_name}  页数: {slide_count}  起始页: {start_slide}")
                cumulative_slide += slide_count

            # 设置文本内容
            text_frame.text = "\n".join(toc_lines)
            
            # 设置字体大小和格式
            for para in text_frame.paragraphs:
                para.font.size = Pt(24)
                para.font.name = "Microsoft YaHei"
                para.space_after = Pt(6)
                para.line_spacing = 1.4

        except Exception as e:
            # 如果创建目录页失败，不影响主流程，只记录错误
            print(f"创建目录页时出错：{e}")

    def _convert_ppts_to_pdfs(self) -> List[Tuple[str, str, bool]]:
        stats: List[Tuple[str, str, bool]] = []

        for item in self.selected_items:
            # 规范化路径，确保使用正确的路径分隔符
            ppt_path = os.path.normpath(item.file_path)
            pdf_path = os.path.normpath(os.path.splitext(ppt_path)[0] + ".pdf")
            existed_before = os.path.exists(pdf_path)
            
            try:
                self._run_vbs_conversion(ppt_path)
            except Exception as exc:
                raise RuntimeError(f"转换 PPT 失败：{item.display_name}\n{exc}") from exc

            # 等待PDF文件生成，最多等待30秒
            max_wait = 30
            wait_interval = 0.5
            waited = 0
            check_count = 0
            while not os.path.exists(pdf_path) and waited < max_wait:
                time.sleep(wait_interval)
                waited += wait_interval
                check_count += 1
                # 每检查10次（约5秒）刷新一次路径（防止路径缓存问题）
                if check_count % 10 == 0:
                    pdf_path = os.path.normpath(os.path.splitext(ppt_path)[0] + ".pdf")

            if not os.path.exists(pdf_path):
                raise RuntimeError(f"未找到转换后的 PDF 文件：{pdf_path}\n请检查PPT文件是否成功转换为PDF。")

            stats.append((item.display_name, pdf_path, existed_before))

        return stats

    def _run_vbs_conversion(self, ppt_path: str):
        """
        调用 VBS 将 PPT 转为 PDF。
        """
        # 确保路径是绝对路径且规范化
        ppt_path = os.path.abspath(os.path.normpath(ppt_path))
        vbs_path = os.path.abspath(os.path.normpath(self.vbs_path))
        
        if not os.path.exists(ppt_path):
            raise RuntimeError(f"PPT文件不存在：{ppt_path}")
        if not os.path.exists(vbs_path):
            raise RuntimeError(f"VBS脚本不存在：{vbs_path}")
        
        cmd = [
            "cscript.exe",
            "//nologo",
            vbs_path,
            ppt_path,
        ]
        completed = subprocess.run(
            cmd, 
            stdout=subprocess.PIPE, 
            stderr=subprocess.PIPE, 
            text=True, 
            shell=False,
            cwd=os.path.dirname(vbs_path)  # 设置工作目录为VBS脚本所在目录
        )
        if completed.returncode != 0:
            error_msg = completed.stderr or completed.stdout or "cscript 返回非零退出码"
            raise RuntimeError(f"VBS转换失败：{error_msg}")

    def _merge_pdfs_with_toc(self, stats: List[Tuple[str, str, bool]], mode_label: str) -> str:
        today_str = datetime.datetime.now().strftime("%Y%m%d")
        base_name = f"{today_str}{mode_label}.pdf"
        output_path = os.path.join(self.folder_path, base_name)

        pdf_infos: List[Tuple[str, str, int]] = []
        for display_name, pdf_path, _existed_before in stats:
            num_pages = self._count_pdf_pages(pdf_path)
            pdf_infos.append((display_name, pdf_path, num_pages))

        toc_pdf_path = self._create_toc_pdf(pdf_infos)
        toc_dir = os.path.dirname(toc_pdf_path)

        try:
            writer = PyPDF2.PdfWriter()

            with open(toc_pdf_path, "rb") as f_toc:
                toc_reader = PyPDF2.PdfReader(f_toc)
                for page in toc_reader.pages:
                    writer.add_page(page)

            for _display_name, pdf_path, _num_pages in pdf_infos:
                with open(pdf_path, "rb") as f_pdf:
                    reader = PyPDF2.PdfReader(f_pdf)
                    for page in reader.pages:
                        writer.add_page(page)

            with open(output_path, "wb") as out_file:
                writer.write(out_file)
        finally:
            shutil.rmtree(toc_dir, ignore_errors=True)
            for _display_name, pdf_path, existed_before in stats:
                if not existed_before and os.path.exists(pdf_path):
                    try:
                        os.remove(pdf_path)
                    except OSError:
                        pass

        return output_path

    def _create_toc_pdf(self, pdf_infos: List[Tuple[str, str, int]]) -> str:
        tmp_dir = tempfile.mkdtemp(prefix="ppt_toc_")
        toc_pdf_path = os.path.join(tmp_dir, "toc.pdf")
        try:
            c = canvas.Canvas(toc_pdf_path, pagesize=A4)
            width, height = A4

            title = "目录"
            title_font = self.font_bold or self.font_regular
            content_font = self.font_regular

            # 适当放大目录标题和正文字号，便于会前快速浏览
            c.setFont(title_font, 36)
            c.drawString(72, height - 72, title)

            c.setFont(content_font, 20)
            y = height - 120
            line_height = 20
            usable_height = height - 72 - 120
            lines_per_page = max(1, int(usable_height // line_height) + 1)

            total_lines = len(pdf_infos)
            toc_page_count = max(1, (total_lines + lines_per_page - 1) // lines_per_page)

            cumulative_page = 0
            toc_lines: List[str] = []

            for idx, (display_name, _pdf_path, num_pages) in enumerate(pdf_infos, start=1):
                start_page = toc_page_count + cumulative_page + 1
                toc_lines.append(f"{idx}. {display_name}  页数: {num_pages}  起始页: {start_page}")
                cumulative_page += num_pages

            for line in toc_lines:
                if y < 72:
                    c.showPage()
                    c.setFont(title_font, 28)
                    c.drawString(72, height - 72, "目录（续）")
                    y = height - 120
                    c.setFont(content_font, 20)
                c.drawString(72, y, line)
                y -= line_height

            c.save()
            return toc_pdf_path
        except Exception:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise

    def _count_pdf_pages(self, pdf_path: str) -> int:
        with open(pdf_path, "rb") as f_pdf:
            reader = PyPDF2.PdfReader(f_pdf)
            return len(reader.pages)


def main():
    if ttkb is not None:
        root = ttkb.Window(themename="cosmo")
    else:
        root = tk.Tk()
    app = PPTMergerApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()

